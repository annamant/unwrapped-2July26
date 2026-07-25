#!/usr/bin/env python3
"""Build Unwrapped brand source deck for Gamma theme import (.pptx)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "Unwrapped-Gamma-Theme-Source.pptx"
LOGO = ROOT.parent / "client" / "public" / "icon-512.png"

# Brand tokens
BG = RGBColor(0xFA, 0xFA, 0xF8)
INK = RGBColor(0x14, 0x12, 0x10)
V = RGBColor(0xE8, 0x34, 0x1C)
MUTED = RGBColor(0x7A, 0x7A, 0x7A)

FONT_DISPLAY = "Playfair Display"
FONT_LABEL = "Space Mono"
FONT_BODY = "DM Sans"


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, *, font, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1 — cover + logo
    s1 = prs.slides.add_slide(blank)
    set_slide_bg(s1, BG)
    if LOGO.exists():
        s1.shapes.add_picture(str(LOGO), Inches(0.8), Inches(0.6), height=Inches(1.1))
    add_textbox(s1, Inches(0.8), Inches(2.0), Inches(11), Inches(1.2), "Unwrapped", font=FONT_DISPLAY, size=54, color=INK, bold=True)
    add_textbox(s1, Inches(0.8), Inches(3.2), Inches(11), Inches(0.8), "Limited local drops from independent shops near you.", font=FONT_BODY, size=22, color=INK)
    add_textbox(s1, Inches(0.8), Inches(4.2), Inches(11), Inches(0.5), "LONDON · shopunwrapped.com", font=FONT_LABEL, size=14, color=MUTED)

    # Slide 2 — typography hierarchy
    s2 = prs.slides.add_slide(blank)
    set_slide_bg(s2, BG)
    add_textbox(s2, Inches(0.8), Inches(0.6), Inches(11), Inches(0.4), "LIMITED LOCAL DROPS", font=FONT_LABEL, size=12, color=V)
    add_textbox(s2, Inches(0.8), Inches(1.2), Inches(11), Inches(1.0), "Headline — Playfair Display", font=FONT_DISPLAY, size=44, color=INK, bold=True)
    add_textbox(s2, Inches(0.8), Inches(2.4), Inches(11), Inches(0.6), "Label — Space Mono", font=FONT_LABEL, size=16, color=INK)
    add_textbox(s2, Inches(0.8), Inches(3.2), Inches(11), Inches(0.8), "Body copy — DM Sans. Reserve in seconds. Collect with QR.", font=FONT_BODY, size=20, color=INK)
    add_textbox(s2, Inches(0.8), Inches(4.2), Inches(11), Inches(0.5), "Accent link colour", font=FONT_BODY, size=18, color=V)

    # Slide 3 — colour swatches (Gamma reads palette from deck)
    s3 = prs.slides.add_slide(blank)
    set_slide_bg(s3, BG)
    add_textbox(s3, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8), "Brand colours", font=FONT_DISPLAY, size=36, color=INK, bold=True)
    swatches = [
        ("Background #FAFAF8", BG, INK),
        ("Ink #141210", INK, BG),
        ("Vermillion #E8341C", V, BG),
        ("Muted text #7A7A7A", MUTED, BG),
    ]
    x = Inches(0.8)
    for label, fill, text in swatches:
        rect = s3.shapes.add_shape(1, x, Inches(2.0), Inches(2.4), Inches(1.6))  # rectangle
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
        rect.line.color.rgb = INK
        add_textbox(s3, x, Inches(3.8), Inches(2.4), Inches(0.8), label, font=FONT_LABEL, size=10, color=INK, align=PP_ALIGN.CENTER)
        x += Inches(2.8)

    # Slide 4 — inverse lockup sample
    s4 = prs.slides.add_slide(blank)
    set_slide_bg(s4, INK)
    add_textbox(s4, Inches(0.8), Inches(2.8), Inches(11), Inches(1.2), "Unwrapped", font=FONT_DISPLAY, size=48, color=BG, bold=True)
    add_textbox(s4, Inches(0.8), Inches(4.0), Inches(11), Inches(0.5), "● LIVE NOW", font=FONT_LABEL, size=14, color=V)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
